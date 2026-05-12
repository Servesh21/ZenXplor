import os, io, logging, threading, time, psutil, urllib.parse
from flask import Blueprint, request, jsonify, current_app, send_file, request as flask_request
from flask_jwt_extended import jwt_required, get_jwt_identity
from models import db, IndexedFile, CloudStorageAccount, User
from elasticsearch import Elasticsearch, helpers
from flask_cors import CORS
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.discovery import build, build_from_document
from google.oauth2.credentials import Credentials
from datetime import datetime, timezone
import dropbox
from dropbox.exceptions import AuthError
from sqlalchemy.orm import scoped_session, sessionmaker
from sqlalchemy import or_
from sqlalchemy.dialects.postgresql import insert
from werkzeug.utils import secure_filename

# Elasticsearch Setup
ELASTICSEARCH_URL = os.getenv("ELASTICSEARCH_URL", "http://localhost:9200")
parsed_url = urllib.parse.urlparse(ELASTICSEARCH_URL)

if "bonsaisearch.net" in ELASTICSEARCH_URL or parsed_url.username:
    # Authenticated cluster (Bonsai/OpenSearch)
    es = Elasticsearch(
        [ELASTICSEARCH_URL],
        use_ssl=True,
        verify_certs=True
    )
    logging.info(f"Connected to authenticated search cluster")
else:
    es = Elasticsearch([ELASTICSEARCH_URL])
    logging.info(f"Connected to search cluster at {ELASTICSEARCH_URL}")

DROPBOX_CLIENT_ID = os.getenv("DROPBOX_CLIENT_ID")
DROPBOX_CLIENT_SECRET = os.getenv("DROPBOX_CLIENT_SECRET")
search_bp = Blueprint("search", __name__)
CORS(search_bp, supports_credentials=True, origins=["http://localhost:5173"])

logging.basicConfig(level=logging.INFO)
indexing_status = {}
from concurrent.futures import ThreadPoolExecutor
executor = ThreadPoolExecutor(max_workers=5)


AUTO_SYNC_INTERVAL_LOCAL = 30
AUTO_SYNC_INTERVAL=30

EXCLUDE_DIRS = {"AppData","node_modules", ".git", ".Trash", "System Volume Information",".venv",".gradle", "Library", ".cache", ".config", ".idea", ".vscode"}
EXCLUDE_FILES = {".DS_Store", "thumbs.db"}

indexing_status = {}

def update_progress(user_id, source, account_id, status, processed, total=None):
    user_id_str = str(user_id)
    if user_id_str not in indexing_status:
        indexing_status[user_id_str] = {}
    if source not in indexing_status[user_id_str]:
        indexing_status[user_id_str][source] = {}
    
    current = indexing_status[user_id_str][source].get(str(account_id), {})
    indexing_status[user_id_str][source][str(account_id)] = {
        "status": status,
        "processed": processed,
        "total": total if total is not None else current.get("total")
    }

auto_sync_started = False  # Global flag


from apscheduler.schedulers.background import BackgroundScheduler
from datetime import datetime

scheduler = BackgroundScheduler()

def start_auto_sync_threads(app):
    """Start background scheduler for auto-syncing local storage, Google Drive, and Dropbox."""
    global scheduler

    if scheduler.running:
        print("🔄 Scheduler already running...")
        return

    print("🚀 Auto-sync scheduler starting...")
    
    # Run immediately, then every 10 minutes
    now = datetime.now()
    scheduler.add_job(func=auto_index_local_storage, args=[app], trigger="interval", minutes=10, id="local_sync", replace_existing=True, next_run_time=now)
    scheduler.add_job(func=auto_index_google_drive, args=[app], trigger="interval", minutes=10, id="gdrive_sync", replace_existing=True, next_run_time=now)
    scheduler.add_job(func=auto_index_dropbox, args=[app], trigger="interval", minutes=10, id="dropbox_sync", replace_existing=True, next_run_time=now)
    
    scheduler.start()
    print("📂 Local storage, ☁️ Google Drive, and 📦 Dropbox auto-sync scheduled (every 10m).")



def get_available_drives():
    """Return all available drives."""
    if os.name == 'nt':  # Windows
        
        return [dp.device for dp in psutil.disk_partitions()]
    else:  # Linux/macOS
        return ["/"]  # Root directory
    
def is_valid_file(file_path):
    """Check if the file should be indexed."""
    file_name = os.path.basename(file_path).lower()
    return (
        file_name not in EXCLUDE_FILES
    )

def is_valid_dir(dir_path):
    """Check if the directory should be indexed."""
    dir_name = os.path.basename(dir_path)
    return dir_name not in EXCLUDE_DIRS

def get_scan_roots():
    """Get the roots to scan for local files. On Windows, scan all drives."""
    if os.name == 'nt':
        return get_available_drives()
    return ["/"]

def auto_index_local_storage(app):
    """Check for new files and index only new ones."""
    global executor
    print("🚀 Local storage auto-indexing starting...")
    
    with app.app_context():
        try:
            users = db.session.query(User.id).all()
            roots = get_scan_roots()
            
            for (user_id,) in users:
                executor.submit(index_local_roots, user_id, roots, app)
        except Exception as e:
            print(f"❌ Error in auto-indexing: {str(e)}")

def index_local_roots(user_id, roots, app):
    """Scan multiple roots and index new files with progress tracking."""
    with app.app_context():
        session = db.session
        all_files_to_index = []
        
        try:
            update_progress(user_id, "local", "default", "fetching", 0, None)
            
            for root in roots:
                print(f"🔍 Counting files in root: {root}")
                for dirpath, dirnames, filenames in os.walk(root):
                    # Exclude directories
                    dirnames[:] = [d for d in dirnames if is_valid_dir(os.path.join(dirpath, d))]
                    
                    for f in filenames:
                        fp = os.path.join(dirpath, f)
                        if is_valid_file(fp):
                            all_files_to_index.append(fp)
            
            total_files = len(all_files_to_index)
            update_progress(user_id, "local", "default", "indexing", 0, total_files)
            
            indexed_items = []
            for idx, file_path in enumerate(all_files_to_index):
                file = os.path.basename(file_path)

                # Get file stats
                file_size = None
                file_mtime = None
                try:
                    file_stat = os.stat(file_path)
                    file_size = file_stat.st_size
                    file_mtime = file_stat.st_mtime
                except OSError:
                    continue

                file_type = file.split('.')[-1] if '.' in file else 'unknown'
                last_mod_dt = datetime.fromtimestamp(file_mtime, timezone.utc) if file_mtime else None

                # Upsert
                stmt = insert(IndexedFile).values(
                    user_id=user_id,
                    filename=file,
                    filepath=file_path,
                    filetype=file_type,
                    is_folder=False,
                    filesize=file_size,
                    last_modified=last_mod_dt,
                    storage_type="local"
                ).on_conflict_do_update(
                    index_elements=["filepath"],
                    set_={
                        "filename": file,
                        "filetype": file_type,
                        "filesize": file_size,
                        "last_modified": last_mod_dt
                    }
                )
                session.execute(stmt)

                indexed_items.append({
                    "id": f"{user_id}_{file_path}",
                    "user_id": user_id,
                    "filename": file,
                    "filepath": file_path,
                    "filetype": file_type,
                    "is_folder": False,
                    "filesize": file_size,
                    "last_modified": file_mtime
                })
                
                if idx > 0 and idx % 100 == 0:
                    update_progress(user_id, "local", "default", "indexing", idx, total_files)

            session.commit()
            
            if indexed_items and es:
                helpers.bulk(es, [{"_index": "file_index", "_id": item["id"], "_source": item} for item in indexed_items])
                
            update_progress(user_id, "local", "default", "completed", total_files, total_files)
            print(f"✅ Local sync complete for user {user_id}. Indexed {len(indexed_items)} files.")

        except Exception as e:
            logging.error(f"❌ Local indexing error: {str(e)}")
            session.rollback()
            update_progress(user_id, "local", "default", "error", 0, None)
        finally:
            session.close()


def auto_index_google_drive(app):
    """Index new files from all users' Google Drive accounts."""
    
    with app.app_context():  # ✅ Ensures app context for database access
        print("🚀 Google Drive Auto-Indexing Started")  # Debug log

        users = db.session.query(IndexedFile.user_id).distinct().all()
        print(f"🔄 Syncing Google Drive for {len(users)} users...")  

        for (user_id,) in users:
            try:
                # ✅ Fetch all linked Google Drive accounts for this user
                account_ids = db.session.query(CloudStorageAccount.id).filter_by(
                    user_id=user_id, provider="Google Drive"
                ).all()
                
                for (account_id,) in account_ids:
                    print(f"🔄 Syncing Google Drive for User {user_id}, Account {account_id}")
                    sync_google_drive(account_id, user_id)  # ✅ Call your function

            except Exception as e:
                logging.error(f"❌ Google Drive indexing error for User {user_id}: {str(e)}")

        db.session.remove()  # ✅ Prevent memory leaks


def auto_index_dropbox(app):
    """Index new files from all users' Dropbox accounts."""
    
    with app.app_context():  # ✅ Ensure Flask app context
        print("🚀 Dropbox Auto-Indexing Started")  # Debug log
        
        users = db.session.query(IndexedFile.user_id).distinct().all()
        print(f"🔄 Syncing Dropbox for {len(users)} users...")  

        for (user_id,) in users:
            try:
                # ✅ Fetch all linked Dropbox accounts for this user
                account_ids = db.session.query(CloudStorageAccount.id).filter_by(
                    user_id=user_id, provider="Dropbox"
                ).all()
                
                for (account_id,) in account_ids:
                    print(f"🔄 Syncing Dropbox for User {user_id}, Account {account_id}")
                    sync_dropbox(account_id, user_id)  # ✅ Call your function

            except Exception as e:
                logging.error(f"❌ Dropbox indexing error for User {user_id}: {str(e)}")

        db.session.remove()  # ✅ Prevent memory leaks

def get_dropbox_access_token(account_id):
    """Fetch the access token for a specific Dropbox account."""
    account = CloudStorageAccount.query.filter_by(id=account_id, provider="Dropbox").first()
    return account.access_token if account else None

def fetch_dropbox_files(dbx, path=""):
    """Recursively fetch files from Dropbox."""
    try:
        result = dbx.files_list_folder(path, recursive=True)
        files = []

        while True:
            for entry in result.entries:
                if isinstance(entry, dropbox.files.FileMetadata):
                    files.append(entry)
            if not result.has_more:
                break
            result = dbx.files_list_folder_continue(result.cursor)

        return files
    except Exception as e:
        logging.error(f"Dropbox API Error: {str(e)}")
        return []
    
def get_dropbox_file_type(file_name):
    """Return the file extension as the file type."""
    if file_name.endswith('/'):
        return 'folder'  # This handles folders, which end with '/'
    else:
        return file_name.split('.')[-1]  # Extract the file extension (e.g., 'pdf', 'jpg', etc.)


def sync_dropbox(account_id, user_id):
    """Sync Dropbox files for a specific user."""
    access_token = get_dropbox_access_token(account_id)
    if not access_token:
        logging.error(f"No valid access token found for account {account_id}")
        return

    dbx = dropbox.Dropbox(access_token)

    try:
        dbx.users_get_current_account()
    except AuthError:
        logging.error("Invalid Dropbox access token")
        return

    with current_app.app_context():
        session = scoped_session(sessionmaker(bind=db.engine))

        try:
            update_progress(user_id, "dropbox", account_id, "fetching", 0, None)
            files = fetch_dropbox_files(dbx)
            total_files = len(files)
            update_progress(user_id, "dropbox", account_id, "indexing", 0, total_files)
            update_count = 0

            for idx, file in enumerate(files):
                file_path = file.path_display
                last_modified = file.server_modified

                # Use the helper function to get the file extension as filetype
                file_type = get_dropbox_file_type(file.name)

                stmt = insert(IndexedFile).values(
                    user_id=user_id,
                    account_id=account_id,
                    filename=file.name,
                    filepath=f"dropbox://{file_path}",
                    storage_type="dropbox",
                    filetype=file_type,  # Save file extension as filetype
                    cloud_file_id=file.id,
                    mime_type=file.content_hash,  # This is a unique identifier for the file, not the MIME type
                    last_modified=last_modified,
                    is_favorite=False  # Default to False for new entries
                ).on_conflict_do_update(
                    index_elements=["filepath"],
                    set_={
                        "filename": file.name,
                        "mime_type": file.content_hash,
                        "filetype": file_type,  # Update filetype based on file extension
                        "last_modified": last_modified,
                        "is_favorite": False  # Default to False for updated entries
                    }
                )

                session.execute(stmt)
                update_count += 1
                
                if idx > 0 and idx % 50 == 0:
                    update_progress(user_id, "dropbox", account_id, "indexing", idx, total_files)

            session.commit()
            update_progress(user_id, "dropbox", account_id, "completed", total_files, total_files)
            logging.info(f"✅ Synced {update_count} files (new + updated) from Dropbox (Account {account_id}) for user {user_id}")

        except Exception as e:
            session.rollback()
            update_progress(user_id, "dropbox", account_id, "error", 0, None)
            logging.error(f"Error syncing Dropbox (Account {account_id}): {str(e)}")

        finally:
            session.remove()


def run_with_app_context(app, func, *args):
    """Runs a function inside the Flask app context in a separate thread."""
    with app.app_context():
        func(*args)



def get_available_drives():
    """Return only the current user's directory in C drive."""
    user_home = os.path.expanduser("~")  # Gets C:/Users/<username>
    return [user_home] if os.path.exists(user_home) else []

def sanitize_filepath(path):
    """Sanitize file paths to ensure consistency."""
    return os.path.abspath(path)




# ---------------------------------------------------------------------------
# Text extraction helpers (used by upload-and-index)
# ---------------------------------------------------------------------------

TEXT_EXTENSIONS = {
    "txt", "md", "rst", "csv", "log", "py", "js", "ts", "jsx", "tsx",
    "html", "htm", "css", "json", "xml", "yaml", "yml", "toml", "ini",
    "cfg", "sh", "bat", "c", "cpp", "h", "java", "rb", "go", "rs",
}

MAX_CONTENT_BYTES = 50_000   # 50 KB plain-text read cap
MAX_CONTENT_SNIPPET = 10_000  # characters stored in Elasticsearch


def extract_text_from_file(file_obj, filetype: str) -> str:
    """
    Extract a text snippet from an uploaded file object.
    Returns an empty string if extraction is not supported or fails.
    """
    try:
        if filetype in TEXT_EXTENSIONS:
            raw = file_obj.read(MAX_CONTENT_BYTES)
            return raw.decode("utf-8", errors="replace")[:MAX_CONTENT_SNIPPET]

        if filetype == "pdf":
            import PyPDF2
            reader = PyPDF2.PdfReader(file_obj)
            parts = []
            for page in reader.pages[:30]:
                text = page.extract_text()
                if text:
                    parts.append(text)
                if sum(len(p) for p in parts) >= MAX_CONTENT_SNIPPET:
                    break
            return "\n".join(parts)[:MAX_CONTENT_SNIPPET]

        if filetype in ("docx",):
            import docx as docx_lib
            doc = docx_lib.Document(file_obj)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:MAX_CONTENT_SNIPPET]

        if filetype == "pptx":
            import pptx as pptx_lib
            prs = pptx_lib.Presentation(file_obj)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)[:MAX_CONTENT_SNIPPET]

    except Exception as e:
        logging.warning(f"Text extraction failed (type={filetype}): {e}")

    return ""


def index_files_worker(user_id, base_directory):
    """Index all folders and files recursively in a given drive/directory with prefix search support."""
    from app import app
    from sqlalchemy.dialects.postgresql import insert
    
    with app.app_context():
        session = scoped_session(sessionmaker(bind=db.engine))  # New DB session

        indexing_status[user_id] = "in_progress"

        try:
            batch_size = 2000
            db_batch = []
            es_batch = []

            def flush_batches():
                nonlocal db_batch, es_batch
                if not db_batch:
                    return
                
                # Bulk UPSERT to PostgreSQL
                stmt = insert(IndexedFile).values(db_batch)
                stmt = stmt.on_conflict_do_update(
                    index_elements=["filepath"],
                    set_={
                        "filename": stmt.excluded.filename,
                        "filetype": stmt.excluded.filetype,
                        "is_folder": stmt.excluded.is_folder,
                        "filesize": stmt.excluded.filesize,
                        "last_modified": stmt.excluded.last_modified
                    }
                )
                session.execute(stmt)
                session.commit()
                print(f"Committed {len(db_batch)} files/folders to DB")

                # Bulk index in Elasticsearch
                if es_batch and es and check_elasticsearch():
                    for doc in es_batch:
                        if "_id" in doc:
                            del doc["_id"]
                    helpers.bulk(es, [{"_index": "file_index", "_id": doc["id"], "_source": doc} for doc in es_batch])
                    print(f"Indexed {len(es_batch)} items in Elasticsearch")
                
                db_batch.clear()
                es_batch.clear()

            for root, dirs, files in os.walk(base_directory):
                # Filter out excluded directories
                dirs[:] = [d for d in dirs if d not in EXCLUDE_DIRS and not d.startswith(".")]
                
                # Index Directories
                for folder in dirs:
                    folder_path = sanitize_filepath(os.path.join(root, folder))
                    
                    db_batch.append({
                        "user_id": user_id,
                        "filename": folder,
                        "filepath": folder_path,
                        "is_folder": True,
                        "filetype": "folder",
                        "storage_type": "local",
                        "is_favorite": False
                    })
                    es_batch.append({
                        "id": folder_path,
                        "user_id": user_id,
                        "filename": folder,
                        "filename_ngram": folder.lower(),
                        "filepath": folder_path,
                        "is_folder": True,
                        "filetype": "folder",
                        "storage_type": "local",
                        "is_favorite": False
                    })

                # Index Files
                for file in files:
                    file_path = sanitize_filepath(os.path.join(root, file))
                    if file in EXCLUDE_FILES or file.startswith("."):
                        continue  # Skip hidden/system files

                    # Extract file extension as filetype
                    _, file_extension = os.path.splitext(file)
                    file_extension = file_extension.lower().strip('.')

                    # Get file stats for size and modification time
                    file_size = None
                    file_mtime = None
                    try:
                        file_stat = os.stat(file_path)
                        file_size = file_stat.st_size
                        file_mtime = file_stat.st_mtime
                    except OSError:
                        pass

                    db_batch.append({
                        "user_id": user_id,
                        "filename": file,
                        "filepath": file_path,
                        "is_folder": False,
                        "filetype": file_extension,
                        "storage_type": "local",
                        "is_favorite": False,
                        "filesize": file_size,
                        "last_modified": datetime.fromtimestamp(file_mtime, timezone.utc) if file_mtime else None
                    })
                    es_batch.append({
                        "id": file_path,
                        "user_id": user_id,
                        "filename": file,
                        "filename_ngram": file.lower(),
                        "filepath": file_path,
                        "is_folder": False,
                        "filetype": file_extension,
                        "storage_type": "local",
                        "is_favorite": False,
                        "filesize": file_size,
                        "last_modified": file_mtime
                    })

                    # Flush if batch size reached
                    if len(db_batch) >= batch_size:
                        flush_batches()

            # Flush remaining
            flush_batches()

        except Exception as e:
            logging.error(f"Error during indexing: {str(e)}")

        finally:
            session.remove()
            indexing_status[user_id] = "completed"

from sqlalchemy.dialects.postgresql import insert


def get_access_token(account_id):
    """Fetch the access token for a specific Google Drive account."""
    account = CloudStorageAccount.query.filter_by(id=account_id, provider="Google Drive").first()
    return account.access_token if account else None

def get_file_type_from_mime(mime_type):
    """Map MIME type to a simplified file type (e.g., 'pdf', 'image', 'docx', etc.)"""
    mime_type_mapping = {
        "application/pdf": "pdf",
        "application/py": "py",
        "application/msword": "docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "image/jpeg": "image",
        "image/png": "image",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "spreadsheet",
        "application/vnd.ms-excel": "spreadsheet",
        "application/zip": "archive",
        "audio/mpeg": "audio",
        "video/mp4": "video",
        "application/vnd.google-apps.spreadsheet": "spreadsheet",
        "application/vnd.jgraph.mxfile": "diagram",
        "application/vnd.google-apps.folder": "folder",
        # Add more mappings as needed
    }
    return mime_type_mapping.get(mime_type, "folder")  # Default to 'unknown' if not found


def sync_google_drive(account_id, user_id):
    """Sync only the specified Google Drive account for a user, resolving conflicts properly."""
    access_token = get_access_token(account_id)
    if not access_token:
        logging.error(f"No valid access token found for account {account_id}")
        return

    with current_app.app_context():
        session = scoped_session(sessionmaker(bind=db.engine))

        account = session.get(CloudStorageAccount, account_id)
        if not account:
            logging.error(f"Account {account_id} not found during sync")
            return

        creds = Credentials(
            token=account.access_token,
            refresh_token=account.refresh_token,
            token_uri="https://oauth2.googleapis.com/token",
            client_id=os.getenv("CLIENT_ID"),
            client_secret=os.getenv("CLIENT_SECRET")
        )
        service = build("drive", "v3", credentials=creds)

        try:
            update_progress(user_id, "google_drive", account_id, "fetching", 0, None)
            files = []
            page_token = None

            # Fetch all files using pagination
            while True:
                response = service.files().list(
                    fields="nextPageToken, files(id, name, mimeType, modifiedTime)",
                    pageSize=100,  # Fetch in batches of 100
                    pageToken=page_token
                ).execute()
                batch_files = response.get("files", [])
                files.extend(batch_files)
                update_progress(user_id, "google_drive", account_id, "fetching", len(files), None)
                
                page_token = response.get("nextPageToken")
                if not page_token:
                    break

            new_entries = []
            update_count = 0
            total_files = len(files)
            
            update_progress(user_id, "google_drive", account_id, "indexing", 0, total_files)

            for idx, file in enumerate(files):
                file_id = file["id"]
                last_modified = datetime.strptime(file["modifiedTime"], "%Y-%m-%dT%H:%M:%S.%fZ")
                mime_type = file["mimeType"]

                # Use the worker function to get the simplified file type
                file_type = get_file_type_from_mime(mime_type)

                # Prepare the insert statement with ON CONFLICT handling
                stmt = insert(IndexedFile).values(
                    user_id=user_id,
                    account_id=account_id,
                    filename=file["name"],
                    filepath=f"drive://{file_id}",
                    storage_type="google_drive",
                    filetype=file_type,  # Correct column name
                    cloud_file_id=file_id,
                    mime_type=mime_type,
                    last_modified=last_modified,
                    is_favorite=False  # Default to False for new entries
                ).on_conflict_do_update(
                    index_elements=["filepath"],  # This ensures the filepath uniqueness constraint is respected
                    set_={
                        "filename": file["name"],
                        "mime_type": mime_type,
                        "filetype": file_type,  # Correct column name
                        "last_modified": last_modified,
                        "is_favorite": False  # Default to False for updated entries
                    }
                )
                
                session.execute(stmt)
                update_count += 1
                
                # Update progress every 50 files
                if idx > 0 and idx % 50 == 0:
                    update_progress(user_id, "google_drive", account_id, "indexing", idx, total_files)

            session.commit()
            update_progress(user_id, "google_drive", account_id, "completed", total_files, total_files)
            logging.info(f"✅ Synced {update_count} files (new + updated) from Google Drive (Account {account_id}) for user {user_id}")

        except Exception as e:
            session.rollback()
            update_progress(user_id, "google_drive", account_id, "error", 0, None)
            logging.error(f"Error syncing Google Drive (Account {account_id}): {str(e)}")

        finally:
            session.remove()


def sync_gmail_attachments(account_id, user_id):
    print("entered gmail attacjments")
    access_token = get_access_token(account_id)
    if not access_token:
        logging.error(f"No valid access token for Gmail (Account {account_id})")
        return
    with current_app.app_context():
        session = scoped_session(sessionmaker(bind=db.engine))
        account = session.get(CloudStorageAccount, account_id)
        if not account:
            logging.error(f"Account {account_id} not found during Gmail sync")
            return

        creds = Credentials(
            token=account.access_token,
            refresh_token=account.refresh_token,
            token_uri="https://oauth2.googleapis.com/token",
            client_id=os.getenv("CLIENT_ID"),
            client_secret=os.getenv("CLIENT_SECRET")
        )
        service = build("gmail", "v1", credentials=creds)

        try:
            messages = service.users().messages().list(userId="me", q="has:attachment").execute().get("messages", [])

            for msg in messages:
                message = service.users().messages().get(userId="me", id=msg["id"]).execute()
                payload = message.get("payload", {})
                parts = payload.get("parts", [])

                for part in parts:
                    if part.get("filename") and "attachmentId" in part.get("body", {}):
                        filename = part["filename"]
                        filetype = filename.split(".")[-1] if "." in filename else "unknown"
                        attachment_id = part["body"]["attachmentId"]

                        stmt = insert(IndexedFile).values(
                            user_id=user_id,
                            account_id=account_id,
                            filename=filename,
                            filepath=f"gmail://{msg['id']}/{attachment_id}",
                            storage_type="gmail",
                            filetype=filetype,
                            cloud_file_id=attachment_id,
                            mime_type=part.get("mimeType", "application/octet-stream"),
                            last_modified=datetime.utcnow()
                        ).on_conflict_do_update(
                            index_elements=["filepath"],
                            set_={"last_modified": datetime.utcnow()}
                        )

                        db.session.execute(stmt)

            db.session.commit()
            logging.info(f"✅ Synced Gmail attachments for user {user_id}")

        except Exception as e:
            session.rollback()
            logging.error(f"Error syncing Gmail attachments: {str(e)}")
        
        finally:
            session.remove()

def sync_google_photos(account_id, user_id):
    access_token = get_access_token(account_id)
    if not access_token:
        logging.error(f"No valid access token for Google Photos (Account {account_id})")
        return

    with current_app.app_context():
        session = scoped_session(sessionmaker(bind=db.engine))
        account = session.get(CloudStorageAccount, account_id)
        if not account:
            logging.error(f"Account {account_id} not found during Google Photos sync")
            return

        creds = Credentials(
            token=account.access_token,
            refresh_token=account.refresh_token,
            token_uri="https://oauth2.googleapis.com/token",
            client_id=os.getenv("CLIENT_ID"),
            client_secret=os.getenv("CLIENT_SECRET")
        )

        photos_api_discovery_url = "https://photoslibrary.googleapis.com/$discovery/rest?version=v1"
        service = build("photoslibrary", "v1", credentials=creds, discoveryServiceUrl=photos_api_discovery_url)


        try:
            update_progress(user_id, "google_photos", account_id, "fetching", 0, None)
            media_items = []
            next_page_token = None

            while True:
                response = service.mediaItems().list(
                    pageSize=100,
                    pageToken=next_page_token
                ).execute()

                media_items.extend(response.get("mediaItems", []))
                update_progress(user_id, "google_photos", account_id, "fetching", len(media_items), None)
                
                next_page_token = response.get("nextPageToken")
                if not next_page_token:
                    break

            total_items = len(media_items)
            update_progress(user_id, "google_photos", account_id, "indexing", 0, total_items)

            for idx, item in enumerate(media_items):
                filename = item.get("filename")
                mime_type = item.get("mimeType", "image/jpeg")
                photo_id = item.get("id")
                filetype = filename.split(".")[-1] if "." in filename else "image"

                stmt = insert(IndexedFile).values(
                    user_id=user_id,
                    account_id=account_id,
                    filename=filename,
                    filepath=f"photos://{photo_id}",
                    storage_type="google_photos",
                    filetype=filetype,
                    cloud_file_id=photo_id,
                    mime_type=mime_type,
                    last_modified=datetime.now(timezone.utc)
                ).on_conflict_do_update(
                    index_elements=["filepath"],
                    set_={"last_modified": datetime.now(timezone.utc)}
                )

                session.execute(stmt)
                
                if idx > 0 and idx % 20 == 0:
                    update_progress(user_id, "google_photos", account_id, "indexing", idx, total_items)

            session.commit()
            update_progress(user_id, "google_photos", account_id, "completed", total_items, total_items)
            logging.info(f"✅ Synced Google Photos for user {user_id}")

        except Exception as e:
            session.rollback()
            update_progress(user_id, "google_photos", account_id, "error", 0, None)
            logging.error(f"Error syncing Google Photos: {str(e)}")
        
        finally:
            session.remove()


@search_bp.route("/sync-agent-files", methods=["POST"])
@jwt_required()
def sync_agent_files():
    """Receive indexed files from the desktop agent and store in PSQL + ES."""
    # Acquire concurrency slot — prevents >3 simultaneous syncs on free tier
    semaphore = current_app.config.get("SYNC_SEMAPHORE")
    acquired = semaphore.acquire(blocking=False) if semaphore else True
    if not acquired:
        return jsonify({"error": "Server busy — retry in a moment"}), 429

    user_id = get_jwt_identity()
    data = request.get_json(silent=True) or {}
    files_batch = data.get("files", [])

    if not files_batch:
        if semaphore and acquired:
            semaphore.release()
        return jsonify({"message": "No files received"}), 200

    session = scoped_session(sessionmaker(bind=db.engine))
    new_entries = []
    es_docs = []

    try:
        for file in files_batch:
            filepath = file.get("filepath", "")[:512]
            filename = file.get("filename", "")[:255]
            filetype = file.get("filetype", "unknown")[:500]
            is_folder = file.get("is_folder", False)
            last_modified = file.get("last_modified")
            
            if last_modified:
                try:
                    last_modified = datetime.fromtimestamp(float(last_modified), timezone.utc)
                except (ValueError, TypeError):
                    last_modified = datetime.now(timezone.utc)
            else:
                last_modified = datetime.now(timezone.utc)

    # Upsert into PostgreSQL
            filesize = file.get("filesize")

            # Upsert into PostgreSQL
            stmt = insert(IndexedFile).values(
                user_id=user_id,
                filename=filename,
                filepath=filepath,
                is_folder=is_folder,
                filetype=filetype,
                storage_type="local",
                is_favorite=False,
                filesize=filesize,
                last_modified=last_modified
            ).on_conflict_do_update(
                index_elements=["filepath"],
                set_={
                    "user_id": user_id,        # keep owner updated on resync
                    "filename": filename,
                    "filetype": filetype,
                    "filesize": filesize,
                    "last_modified": last_modified,
                }
            )
            session.execute(stmt)

            # Prepare Elasticsearch Payload
            es_docs.append({
                "id": filepath,
                "user_id": user_id,
                "filename": filename,
                "filename_ngram": filename.lower(),
                "filepath": filepath,
                "is_folder": is_folder,
                "filetype": filetype,
                "storage_type": "local",
                "is_favorite": False,
                "filesize": filesize,
                "last_modified": last_modified.isoformat() if last_modified else None
            })

        # Commit all records to PostgreSQL
        try:
            session.commit()
            logging.info(f"sync_agent_files: committed {len(files_batch)} records to PostgreSQL for user {user_id}")
        except Exception as db_err:
            session.rollback()
            logging.error(f"sync_agent_files: PostgreSQL commit failed for user {user_id}: {db_err}")
            return jsonify({"error": "Database commit failed", "details": str(db_err)}), 500

        # Push to Elasticsearch
        if es_docs and es and check_elasticsearch():
            try:
                for doc in es_docs:
                    if "_id" in doc:
                        del doc["_id"]
                helpers.bulk(es, [{"_index": "file_index", "_id": doc["id"], "_source": doc} for doc in es_docs])
                logging.info(f"sync_agent_files: indexed {len(es_docs)} docs in Elasticsearch for user {user_id}")
            except Exception as es_err:
                # ES failure is non-fatal — PostgreSQL data is already saved.
                logging.warning(f"sync_agent_files: Elasticsearch bulk failed (non-fatal): {es_err}")
        elif not check_elasticsearch():
            logging.warning("sync_agent_files: Elasticsearch is unavailable — skipping ES indexing. PostgreSQL records were saved.")

        return jsonify({"message": f"Successfully synced {len(files_batch)} files"}), 200

    except Exception as e:
        session.rollback()
        logging.error(f"sync_agent_files: Unexpected error for user {user_id}: {str(e)}")
        error_details = str(e)
        if hasattr(e, 'errors'):
            error_details = str(e.errors)
        return jsonify({"error": "Failed to sync files", "details": error_details}), 500
    finally:
        session.remove()
        if semaphore and acquired:
            semaphore.release()



@search_bp.route("/index-files", methods=["POST"])
@jwt_required()
def index_files():
    """Start indexing all drives and files recursively."""
    user_id = get_jwt_identity()
    
    if not check_elasticsearch():
        return jsonify({"error": "Elasticsearch is not available"}), 500

    drives = get_available_drives()
    indexing_status[user_id] = "starting"

    for drive in drives:
        threading.Thread(target=index_files_worker, args=(user_id, drive)).start()

    return jsonify({"message": f"Indexing started for drives: {drives}"}), 202



@search_bp.route("/upload-and-index", methods=["POST"])
@jwt_required()
def upload_and_index():
    """
    Accept multipart file uploads, extract text content, and index metadata +
    content in Elasticsearch.  Files are NOT stored on the server — only
    metadata and extracted text snippets are persisted so users can search by
    content.  Works on every platform (desktop and mobile) without any
    additional software.
    """
    user_id = get_jwt_identity()
    uploaded_files = request.files.getlist("files")

    if not uploaded_files:
        return jsonify({"error": "No files provided"}), 400

    if len(uploaded_files) > 500:
        return jsonify({"error": "Too many files in a single request (max 500)"}), 400

    session = scoped_session(sessionmaker(bind=db.engine))
    db_rows = []
    es_docs = []
    skipped = 0

    try:
        for f in uploaded_files:
            raw_name = f.filename or ""
            filename = secure_filename(raw_name)
            if not filename:
                skipped += 1
                continue

            _, ext = os.path.splitext(filename)
            filetype = ext.lower().strip(".") or "unknown"

            # Use the original relative path supplied by the browser when
            # available (webkitRelativePath via a same-name form field),
            # otherwise fall back to just the filename.
            relative_path = (request.form.get(f"paths[{raw_name}]") or filename).strip("/")
            filepath = f"upload://{user_id}/{relative_path}"
            if len(filepath) > 512:
                filepath = f"upload://{user_id}/{filename}"

            content_snippet = extract_text_from_file(f.stream, filetype)
            mime_type = f.mimetype or None

            db_rows.append({
                "user_id": user_id,
                "filename": filename[:255],
                "filepath": filepath[:512],
                "is_folder": False,
                "filetype": filetype[:500],
                "storage_type": "local_upload",
                "mime_type": (mime_type[:1024] if mime_type else None),
                "last_modified": datetime.now(timezone.utc),
                "is_favorite": False,
            })

            es_docs.append({
                "id": filepath[:512],
                "user_id": user_id,
                "filename": filename[:255],
                "filename_ngram": filename.lower()[:255],
                "filepath": filepath[:512],
                "is_folder": False,
                "filetype": filetype[:500],
                "storage_type": "local_upload",
                "is_favorite": False,
                "file_content": content_snippet,
            })

        if not db_rows:
            return jsonify({"error": "No valid files to index"}), 400

        for row in db_rows:
            stmt = insert(IndexedFile).values(**row).on_conflict_do_update(
                index_elements=["filepath"],
                set_={
                    "filename": row["filename"],
                    "filetype": row["filetype"],
                    "storage_type": row["storage_type"],
                    "mime_type": row["mime_type"],
                    "last_modified": row["last_modified"],
                },
            )
            session.execute(stmt)

        session.commit()

        if es and check_elasticsearch():
            helpers.bulk(
                es,
                [
                    {"_index": "file_index", "_id": doc["id"], "_source": doc}
                    for doc in es_docs
                ],
            )

        return jsonify({
            "message": "Files uploaded and indexed successfully",
            "indexed_count": len(db_rows),
            "skipped_count": skipped,
        }), 200

    except Exception as e:
        session.rollback()
        logging.error(f"Upload indexing failed: {str(e)}")
        return jsonify({"error": "Failed to upload and index files"}), 500
    finally:
        session.remove()



@search_bp.route("/index-status", methods=["GET"])
@jwt_required()
def get_index_status():
    """Check indexing status for the user."""
    user_id = get_jwt_identity()
    status = indexing_status.get(user_id, "not_started")
    return jsonify({"status": status})

@search_bp.route("/search-files", methods=["GET"])
@jwt_required()
def search_files():
    """Search files with fuzzy matching and pagination from both Elasticsearch and DB."""

    if not check_elasticsearch():
        return jsonify({"error": "Search service unavailable"}), 500

    user_id = get_jwt_identity()
    query = request.args.get("q", "").strip()
    limit = request.args.get("limit", 10, type=int)
    offset = request.args.get("offset", 0, type=int)

    service_filter = request.args.get("service", "").lower()
    filetype_filter = request.args.get("filetype", "").lower()

    # Smart filter params
    modified_after  = request.args.get("modified_after", "").strip()
    modified_before = request.args.get("modified_before", "").strip()
    size_min        = request.args.get("size_min", type=int)
    size_max        = request.args.get("size_max", type=int)

    # Record search query
    if query and offset == 0: # Only record the first page of search
        try:
            from models import SearchHistory
            # Check if this query was recently added (last 1 minute) to avoid duplicates from re-renders
            # Use db.session.query() because SearchHistory.query is shadowed by the 'query' column
            recent_exists = db.session.query(SearchHistory).filter_by(user_id=user_id, query=query).order_by(SearchHistory.timestamp.desc()).first()
            
            should_add = True
            if recent_exists:
                time_diff = datetime.utcnow() - recent_exists.timestamp
                if time_diff.total_seconds() < 60:
                    should_add = False
            
            if should_add:
                new_search = SearchHistory(user_id=user_id, query=query)
                db.session.add(new_search)
                db.session.commit()
        except Exception as e:
            logging.error(f"Failed to record search query: {e}")
            db.session.rollback()

    # If query is empty, we just return recent activity: files only, sorted by modification
    if not query:
        try:
            with current_app.app_context():
                session = scoped_session(sessionmaker(bind=db.engine))
                filters = [IndexedFile.user_id == user_id, IndexedFile.is_folder == False, IndexedFile.last_accessed.isnot(None)]
                
                if service_filter:
                    if service_filter == "local":
                        filters.append(or_(IndexedFile.storage_type == "local", IndexedFile.storage_type == "local_upload"))
                    else:
                        filters.append(IndexedFile.storage_type == service_filter)
                if filetype_filter:
                    filters.append(IndexedFile.filetype == filetype_filter)

                # Recent activity is limited to a fixed set of the most recent items
                RECENT_LIMIT = 24
                recent_files = session.query(IndexedFile).filter(*filters)\
                    .order_by(IndexedFile.last_accessed.desc())\
                    .limit(RECENT_LIMIT).all()
                
                results = [{
                    "id": f.id,
                    "filename": f.filename,
                    "filepath": f.filepath,
                    "storage_type": f.storage_type,
                    "cloud_file_id": f.cloud_file_id,
                    "is_favorite": f.is_favorite
                } for f in recent_files]
                
                session.remove()
                
                # For recent activity, we don't paginate past the fixed set
                return jsonify({
                    "results": results,
                    "total_results": len(results),
                    "has_more": False,
                    "offset": len(results)
                }), 200
        except Exception as e:
            logging.error(f"Recent files fetch error: {e}")
            return jsonify({"results": [], "has_more": False}), 200

    all_results = []
    seen_keys = set()

    def get_unique_key(storage_type, path):
        return f"{storage_type}::{path}"

    # 1️⃣ Elasticsearch results
    try:
        should_clauses = []

        if "*" in query:
            should_clauses.append({
                "wildcard": {
                    "filename": {
                        "value": query,
                        "case_insensitive": True
                    }
                }
            })
        else:
            should_clauses = [
                # 1. Exact match or very close match (High boost)
                {
                    "match": {
                        "filename": {
                            "query": query,
                            "boost": 10
                        }
                    }
                },
                # 2. Prefix match (Medium boost)
                {
                    "prefix": {
                        "filename": {
                            "value": query.lower(),
                            "boost": 5
                        }
                    }
                },
                # 3. N-gram / partial match
                {
                    "match": {
                        "filename_ngram": {
                            "query": query,
                            "boost": 2
                        }
                    }
                },
                # 4. Fuzzy match (Low boost, helps with typos but doesn't override exacts)
                {
                    "match": {
                        "filename": {
                            "query": query,
                            "fuzziness": "AUTO",
                            "boost": 1
                        }
                    }
                },
                # 5. Content match (Lowest boost)
                {
                    "match": {
                        "file_content": {
                            "query": query,
                            "fuzziness": "AUTO",
                            "boost": 0.5
                        }
                    }
                }
            ]

        es_query = {
            "query": {
                "bool": {
                    "should": should_clauses,
                    "minimum_should_match": 1,
                    "filter": [
                        {"term": {"user_id": user_id}}
                    ]
                }
            },
            "size": 100  # Get max possible, pagination happens after merge
        }

        if service_filter:
            if service_filter == "local":
                es_query["query"]["bool"]["filter"].append({
                    "terms": {"storage_type": ["local", "local_upload"]}
                })
            else:
                es_query["query"]["bool"]["filter"].append({
                    "term": {"storage_type": service_filter}
                })

        if filetype_filter:
            es_query["query"]["bool"]["filter"].append({
                "term": {"filetype": filetype_filter}
            })

        # Date range filter
        if modified_after or modified_before:
            date_range = {}
            if modified_after:
                date_range["gte"] = modified_after
            if modified_before:
                date_range["lte"] = modified_before
            es_query["query"]["bool"]["filter"].append({
                "range": {"last_modified": date_range}
            })

        # File size filter
        if size_min is not None or size_max is not None:
            size_range = {}
            if size_min is not None:
                size_range["gte"] = size_min
            if size_max is not None:
                size_range["lte"] = size_max
            es_query["query"]["bool"]["filter"].append({
                "range": {"filesize": size_range}
            })

        es_results = es.search(index="file_index", body=es_query)
        for hit in es_results["hits"]["hits"]:
            doc = hit["_source"]
            path = doc.get("filepath") or doc.get("cloud_file_id")
            key = get_unique_key(doc.get("storage_type"), path)
            if key not in seen_keys:
                all_results.append(doc)
                seen_keys.add(key)
            

    except Exception as e:
        logging.error(f"Elasticsearch error: {str(e)}")

    # 2️⃣ DB fallback (always run alongside ES)
    try:
        with current_app.app_context():
            session = scoped_session(sessionmaker(bind=db.engine))

            filters = [IndexedFile.user_id == user_id]

            if query:
                filters.append(IndexedFile.filename.ilike(f"%{query.strip('*')}%"))
            if service_filter:
                if service_filter == "local":
                    filters.append(or_(
                        IndexedFile.storage_type == "local",
                        IndexedFile.storage_type == "local_upload",
                    ))
                else:
                    filters.append(IndexedFile.storage_type == service_filter)
            if filetype_filter:
                filters.append(IndexedFile.filetype == filetype_filter)

            # Smart filters — date range
            if modified_after:
                try:
                    dt = datetime.strptime(modified_after, "%Y-%m-%d")
                    filters.append(IndexedFile.last_modified >= dt)
                except ValueError:
                    pass

            if modified_before:
                try:
                    dt = datetime.strptime(modified_before, "%Y-%m-%d")
                    filters.append(IndexedFile.last_modified <= dt)
                except ValueError:
                    pass

            # Smart filters — file size
            if size_min is not None:
                filters.append(IndexedFile.filesize >= size_min)

            if size_max is not None:
                filters.append(IndexedFile.filesize <= size_max)

            db_files = session.query(IndexedFile).filter(*filters).all()
            session.remove()

            for file in db_files:
                path = file.filepath or file.cloud_file_id
                key = get_unique_key(file.storage_type, path)
                if key not in seen_keys:
                    all_results.append({
                        "filename": file.filename,
                        "cloud_file_id": file.cloud_file_id,
                        "storage_type": file.storage_type,
                        "filepath": file.filepath,
                        "is_favorite": file.is_favorite,
                        "filetype": file.filetype,
                        "filesize": file.filesize,
                        "last_modified": str(file.last_modified) if file.last_modified else None,
                        "mime_type": file.mime_type,
                    })
                    seen_keys.add(key)
    except Exception as e:
        logging.error(f"DB fallback error: {str(e)}")

    # 3️⃣ Pagination after merging
    total_results = len(all_results)
    paginated = all_results[offset:offset + limit]
    has_more = offset + limit < total_results

    return jsonify({
        "results": paginated,
        "total_results": total_results,
        "offset": offset + len(paginated),
        "limit": limit,
        "has_more": has_more
    }), 200


@search_bp.route("/open-file", methods=["POST"])
@jwt_required()
def open_file():
    """Open file location in File Explorer, Finder, or File Manager, or return cloud file URLs."""
    user_id = get_jwt_identity()
    data = request.get_json()
    file_id = data.get("id")
    file_path = data.get("filepath")

    # 🔍 Debugging Logs
    print(f"User {user_id} is requesting to open: id={file_id}, path={file_path}")

    if not file_id and not file_path:
        return jsonify({"error": "File ID or path is required"}), 400

    # Check if the file is indexed and belongs to the authenticated user
    if file_id:
        file_record = IndexedFile.query.filter_by(id=file_id, user_id=user_id).first()
    else:
        file_record = IndexedFile.query.filter_by(filepath=file_path, user_id=user_id).first()
        
    if not file_record:
        print("❌ Unauthorized access attempt or file not found")  # Debugging log
        return jsonify({"error": "File not found or unauthorized access"}), 403

    storage_type = file_record.storage_type  # "local", "google_drive", "dropbox"
    cloud_file_id = file_record.cloud_file_id  # Used for Google Drive & Dropbox

    try:
        if storage_type == "local":
            return jsonify({"error": "Use the ZenXplor Desktop Agent to open local files"}), 400

        elif storage_type == "google_drive":
            if not cloud_file_id:
                return jsonify({"error": "Cloud file ID missing"}), 400

            drive_url = f"https://drive.google.com/file/d/{cloud_file_id}/view?usp=drivesdk"
            return jsonify({"url": drive_url})

        elif storage_type == "dropbox":
            # Use Dropbox preview link or create a shared link
            account = CloudStorageAccount.query.filter_by(user_id=user_id, provider="Dropbox").first()
            if not account or not account.access_token:
                return jsonify({"error": "No linked Dropbox account"}), 400
            
            try:
                dbx = dropbox.Dropbox(account.access_token)
                db_path = file_record.filepath.replace("dropbox://", "")
                
                # Try to get existing shared link first to avoid creating many links
                shared_links = dbx.sharing_list_shared_links(path=db_path, direct_only=True).links
                if shared_links:
                    dropbox_url = shared_links[0].url
                else:
                    # Create a new shared link
                    res = dbx.sharing_create_shared_link_with_settings(db_path)
                    dropbox_url = res.url
                
                return jsonify({"url": dropbox_url})
            except Exception as e:
                # Fallback to simple preview URL if sharing fails
                db_path = file_record.filepath.replace("dropbox://", "")
                dropbox_url = f"https://www.dropbox.com/preview{db_path}"
                return jsonify({"url": dropbox_url})

        else:
            return jsonify({"error": "Unsupported storage type"}), 400

    except Exception as e:
        print(f"❌ Failed to open file: {e}")  # Debugging log
        return jsonify({"error": f"Failed to open file: {str(e)}"}), 500
    
@search_bp.route("/sync-cloud-storage", methods=["POST"])
@jwt_required()
def sync_google_drive_account():
    """Sync a specific Google Drive account for the user in the background."""
    user_id = get_jwt_identity()
    data = request.json
    account_id = data.get("account_id")

    if not account_id:
        return jsonify({"error": "Account ID is required"}), 400

    try:
        executor.submit(sync_google_drive, account_id, user_id)
        return jsonify({"message": f"Google Drive sync started for account {account_id}"}), 200
    except Exception as e:
        logging.error(f"Error starting Google Drive sync (Account {account_id}): {str(e)}")
        return jsonify({"error": "Failed to start Google Drive sync"}), 500

@search_bp.route("/gmail/sync", methods=["POST"])
@jwt_required()
def sync_gmail():
    """Sync Gmail attachments for the user in the background."""
    user_id = get_jwt_identity()
    data = request.json
    account_id = data.get("account_id")

    if not account_id:
        return jsonify({"error": "Account ID is required"}), 400

    try:
        executor.submit(sync_gmail_attachments, account_id, user_id)
        return jsonify({"message": f"Gmail sync started for account {account_id}"}), 200
    except Exception as e:
        logging.error(f"Error starting Gmail sync (Account {account_id}): {str(e)}")
        return jsonify({"error": "Failed to start Gmail sync"}), 500

@search_bp.route("/photos/sync", methods=["POST"])
@jwt_required()
def sync_gphotos():
    """Sync Google Photos for the user."""
    user_id = get_jwt_identity()
    data = request.json
    account_id = data.get("account_id")

    if not account_id:
        return jsonify({"error": "Account ID is required"}), 400

    try:
        executor.submit(sync_google_photos, account_id, user_id)
        return jsonify({"message": f"Google Photos sync started for account {account_id}"}), 200
    except Exception as e:
        logging.error(f"Error starting Google Photos sync (Account {account_id}): {str(e)}")
        return jsonify({"error": "Failed to start Google Photos sync"}), 500

@search_bp.route("/indexing-progress", methods=["GET"])
@jwt_required()
def get_indexing_progress():
    """Return the real-time indexing progress for the authenticated user."""
    user_id = str(get_jwt_identity())
    return jsonify(indexing_status.get(user_id, {})), 200

def check_elasticsearch():
    """Check if Elasticsearch is reachable."""
    if not es.ping():
        logging.error("Elasticsearch server is not reachable")
        return False
    return True

@search_bp.route("/sync-dropbox", methods=["POST"])
@jwt_required()
def sync_dropbox_account():
    """Sync a specific Dropbox account for the user."""
    user_id = get_jwt_identity()
    data = request.json
    account_id = data.get("account_id")

    if not account_id:
        return jsonify({"error": "Account ID is required"}), 400

    try:
        executor.submit(sync_dropbox, account_id, user_id)
        return jsonify({"message": f"Dropbox sync started for account {account_id}"}), 200
    except Exception as e:
        logging.error(f"Error starting Dropbox sync (Account {account_id}): {str(e)}")
        return jsonify({"error": "Failed to start Dropbox sync"}), 500

@search_bp.route("/download-file", methods=["GET"])
@jwt_required()
def download_file():
    """Allow authenticated users to download a file, including Google Drive files."""
    user_id = get_jwt_identity()
    file_id = flask_request.args.get("id")
    file_path = flask_request.args.get("filepath")

    if not file_id and not file_path:
        return jsonify({"error": "File ID or path is required"}), 400

    if file_id:
        file_record = IndexedFile.query.filter_by(id=file_id, user_id=user_id).first()
    else:
        file_record = IndexedFile.query.filter_by(filepath=file_path, user_id=user_id).first()

    if not file_record:
        return jsonify({"error": "File not found or unauthorized access"}), 403

    # If the file is a local file, return it
    if file_record.storage_type in ["local", "local_upload"]:
        # Handle local_upload paths mapping to the actual filesystem
        actual_path = file_path
        if file_record.storage_type == "local_upload" and file_path.startswith("upload://"):
            # Path looks like upload://user_id/relative_path
            # But the actual file should be in the uploads directory
            # For now, we assume file_path is already set correctly or handled by the indexer
            pass

        if not os.path.exists(actual_path):
            return jsonify({"error": f"File not found: {actual_path}"}), 404
            
        is_inline = flask_request.args.get("inline") == "1"
        return send_file(actual_path, as_attachment=not is_inline)

    # For Dropbox
    if file_record.storage_type == "dropbox":
        account = None
        if file_record.account_id:
            account = CloudStorageAccount.query.get(file_record.account_id)
        if not account:
            account = CloudStorageAccount.query.filter_by(user_id=user_id, provider="Dropbox").first()
            
        if not account or not account.access_token:
            return jsonify({"error": "No linked Dropbox account"}), 400
        
        try:
            dbx = dropbox.Dropbox(account.access_token)
            db_path = file_record.filepath.replace("dropbox://", "")
            metadata, res = dbx.files_download(db_path)
            file_stream = io.BytesIO(res.content)
            return send_file(
                file_stream,
                as_attachment=True,
                download_name=file_record.filename,
                mimetype=file_record.mime_type or "application/octet-stream"
            )
        except Exception as e:
            logging.error(f"Failed to download file from Dropbox: {str(e)}")
            return jsonify({"error": f"Failed to download file from Dropbox: {str(e)}"}), 500

    # For Google Drive files, download from Google Drive
    if file_record.storage_type == "google_drive":
        account = None
        if file_record.account_id:
            account = CloudStorageAccount.query.get(file_record.account_id)
        if not account:
            account = CloudStorageAccount.query.filter_by(user_id=user_id, provider="Google Drive").first()

        if not account:
            return jsonify({"error": "No linked Google Drive account"}), 400

        creds = Credentials(
            token=account.access_token,
            refresh_token=account.refresh_token,
            token_uri="https://oauth2.googleapis.com/token",
            client_id=os.getenv("CLIENT_ID"),
            client_secret=os.getenv("CLIENT_SECRET")
        )
        drive_service = build("drive", "v3", credentials=creds)
        file_id = file_record.cloud_file_id
        mime_type = file_record.mime_type

        try:
            if mime_type and mime_type.startswith("application/vnd.google-apps."):
                # Handle Google Docs, Sheets, etc. by exporting them
                export_mime = "application/pdf"
                if "document" in mime_type:
                    export_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    ext = ".docx"
                elif "spreadsheet" in mime_type:
                    export_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    ext = ".xlsx"
                elif "presentation" in mime_type:
                    export_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    ext = ".pptx"
                else:
                    ext = ".pdf"
                
                drive_request = drive_service.files().export_media(fileId=file_id, mimeType=export_mime)
                download_name = file_record.filename
                if not download_name.lower().endswith(ext):
                    download_name += ext
            else:
                drive_request = drive_service.files().get_media(fileId=file_id)
                download_name = file_record.filename

            file_stream = io.BytesIO()
            downloader = MediaIoBaseDownload(file_stream, drive_request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            file_stream.seek(0)
            return send_file(
                file_stream,
                as_attachment=True,
                download_name=download_name,
                mimetype=mime_type or "application/octet-stream"
            )
        except Exception as e:
            msg = str(e)
            # Specifically check for permission errors (scope issues)
            if "appNotAuthorizedToFile" in msg or "403" in msg:
                return jsonify({
                    "error": "Google Drive permission denied. Please re-connect your Google account in Settings and ensure you grant 'Full Drive' or 'Read-only' access to allow downloads.",
                    "code": "permission_denied"
                }), 403
                
            logging.error(f"Failed to download file from Google Drive: {msg}")
            return jsonify({"error": f"Failed to download file from Google Drive: {msg}"}), 500
    


@search_bp.route("/preview", methods=["GET"])
@jwt_required()
def preview_file():
    """
    Return preview data for a file.
    For images: return base64 encoded thumbnail (max 800px wide)
    For text files: return first 5000 characters of content
    For PDFs: return a Google Docs viewer URL
    For cloud files: return the appropriate viewer/embed URL
    For everything else: return metadata only with preview_type = "none"
    """
    user_id = get_jwt_identity()
    filepath = request.args.get("filepath")

    if not filepath:
        return jsonify({"error": "filepath required"}), 400

    file_record = IndexedFile.query.filter_by(
        filepath=filepath, user_id=user_id
    ).first()
    if not file_record:
        return jsonify({"error": "File not found"}), 404

    storage_type = file_record.storage_type
    filename = file_record.filename
    filetype = (file_record.filetype or "").lower()

    # --- Image preview (local files only) ---
    IMAGE_TYPES = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "svg"}
    TEXT_TYPES  = {"txt", "md", "csv", "json", "xml", "yaml", "yml",
                   "py", "js", "ts", "jsx", "tsx", "html", "css",
                   "sh", "bat", "log", "env", "toml", "ini", "cfg"}

    if storage_type in ("local", "local_browser") and filetype in IMAGE_TYPES:
        try:
            from PIL import Image
            import base64
            with Image.open(filepath) as img:
                img.thumbnail((800, 800))
                buf = io.BytesIO()
                fmt = "PNG" if filetype == "png" else "JPEG"
                img.save(buf, format=fmt)
                b64 = base64.b64encode(buf.getvalue()).decode()
                return jsonify({
                    "preview_type": "image",
                    "data": f"data:image/{filetype};base64,{b64}",
                    "filename": filename,
                    "filetype": filetype,
                    "filesize": file_record.filesize,
                    "last_modified": str(file_record.last_modified),
                    "storage_type": storage_type,
                })
        except Exception as e:
            logging.error(f"Image preview failed: {e}")
            return jsonify({"preview_type": "none", "filename": filename})

    # --- Text preview (local files only) ---
    if storage_type in ("local", "local_browser") and filetype in TEXT_TYPES:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(5000)
            return jsonify({
                "preview_type": "text",
                "content": content,
                "filename": filename,
                "filetype": filetype,
                "filesize": file_record.filesize,
                "last_modified": str(file_record.last_modified),
                "storage_type": storage_type,
            })
        except Exception as e:
            logging.error(f"Text preview failed: {e}")

    # --- PDF preview (local) ---
    if storage_type in ("local", "local_browser") and filetype == "pdf":
        encoded_path = urllib.parse.quote(filepath)
        # Point to our own download-file endpoint which will serve the PDF over HTTP
        # Use inline=1 to ensure the browser displays it instead of downloading
        viewer_url = f"{request.host_url.rstrip('/')}/search/download-file?filepath={encoded_path}&inline=1"
        return jsonify({
            "preview_type": "pdf",
            "viewer_url": viewer_url,
            "download_url": f"/search/download-file?filepath={filepath}",
            "filename": filename,
            "filetype": filetype,
            "filesize": file_record.filesize,
            "last_modified": str(file_record.last_modified),
            "storage_type": storage_type,
        })

    # --- Google Drive ---
    if storage_type == "google_drive" and file_record.cloud_file_id:
        fid = file_record.cloud_file_id
        viewer_url = f"https://drive.google.com/file/d/{fid}/preview"
        return jsonify({
            "preview_type": "iframe",
            "viewer_url": viewer_url,
            "filename": filename,
            "filetype": filetype,
            "storage_type": storage_type,
            "last_modified": str(file_record.last_modified),
        })

    # --- Dropbox ---
    if storage_type == "dropbox" and file_record.cloud_file_id:
        return jsonify({
            "preview_type": "iframe",
            "viewer_url": f"https://www.dropbox.com/preview{file_record.filepath.replace('dropbox:/', '')}?role=personal",
            "filename": filename,
            "filetype": filetype,
            "storage_type": storage_type,
            "last_modified": str(file_record.last_modified),
        })

    # --- Gmail attachment ---
    if storage_type == "gmail":
        return jsonify({
            "preview_type": "none",
            "filename": filename,
            "filetype": filetype,
            "storage_type": storage_type,
            "last_modified": str(file_record.last_modified),
            "message": "Open in Gmail to preview this attachment.",
        })

    # --- Fallback ---
    return jsonify({
        "preview_type": "none",
        "filename": filename,
        "filetype": filetype,
        "filesize": getattr(file_record, "filesize", None),
        "last_modified": str(file_record.last_modified),
        "storage_type": storage_type,
    })


@search_bp.route('/<string:file_path>/favorite', methods=['POST'])
@jwt_required()
def toggle_favorite(file_path):
    user_id = get_jwt_identity()
    file = IndexedFile.query.filter_by(filepath=file_path, user_id=user_id).first()

    
    if not file:
        return jsonify({"error": "File not found"}), 404

    # data = request.get_json()
    # is_favorite = data.get('is_favorite')



    file.is_favorite = not file.is_favorite
    db.session.commit()

    return jsonify({"message": "Favorite status updated", "file": file.to_dict()})


@search_bp.route('/recent-searches', methods=['GET'])
@jwt_required()
def get_recent_searches():
    user_id = get_jwt_identity()
    try:
        from models import SearchHistory
        # Get last 50 queries, then deduplicate in python
        # Use db.session.query() because SearchHistory.query is shadowed by the 'query' column
        history = db.session.query(SearchHistory).filter_by(user_id=user_id).order_by(SearchHistory.timestamp.desc()).limit(50).all()
        
        seen = set()
        unique_queries = []
        for h in history:
            if h.query not in seen:
                unique_queries.append(h.query)
                seen.add(h.query)
            if len(unique_queries) >= 10: # Return top 10 unique recent searches
                break
                
        return jsonify({"recent_searches": unique_queries}), 200
    except Exception as e:
        logging.error(f"Failed to fetch recent searches: {e}")
        return jsonify({"error": "Failed to fetch recent searches"}), 500

@search_bp.route('/file/access', methods=['POST'])
@jwt_required()
def log_file_access():
    user_id = get_jwt_identity()
    data = request.get_json()
    
    file = None
    if data.get('filepath'):
        file = IndexedFile.query.filter_by(filepath=data['filepath'], user_id=user_id).first()
    elif data.get('cloud_file_id'):
        file = IndexedFile.query.filter_by(cloud_file_id=data['cloud_file_id'], user_id=user_id).first()
    elif data.get('id'):
        file = IndexedFile.query.filter_by(id=data['id'], user_id=user_id).first()
        
    if not file:
        return jsonify({"error": "File not found"}), 404
        
    file.last_accessed = datetime.utcnow()
    db.session.commit()
    
    return jsonify({"message": "Access logged"}), 200


@search_bp.route('/storage/stats', methods=['GET'])
@jwt_required()
def get_storage_stats():
    user_id = get_jwt_identity()
    
    stats = {
        "local": {"connected": True, "count": 0, "status": "Connected"},
        "google_drive": {"connected": False, "count": 0, "status": "Not Connected"},
        "dropbox": {"connected": False, "count": 0, "status": "Not Connected"},
        "gmail": {"connected": False, "count": 0, "status": "Not Connected"}
    }
    
    try:
        # Get file counts per storage_type
        counts = db.session.query(
            IndexedFile.storage_type, 
            db.func.count(IndexedFile.id)
        ).filter(
            IndexedFile.user_id == user_id,
            IndexedFile.is_folder == False
        ).group_by(IndexedFile.storage_type).all()
        
        local_count = sum([count for storage_type, count in counts if storage_type in ('local', 'local_upload')])
        stats["local"]["count"] = local_count
        
        gdrive_count = sum([count for storage_type, count in counts if storage_type == 'google_drive'])
        stats["google_drive"]["count"] = gdrive_count
        
        dropbox_count = sum([count for storage_type, count in counts if storage_type == 'dropbox'])
        stats["dropbox"]["count"] = dropbox_count
        
        # Update connection status based on linked accounts
        accounts = CloudStorageAccount.query.filter_by(user_id=user_id).all()
        
        for account in accounts:
            provider = account.provider.lower()
            if 'google' in provider or 'drive' in provider:
                stats["google_drive"]["connected"] = True
                stats["google_drive"]["status"] = "Connected"
            elif 'dropbox' in provider:
                stats["dropbox"]["connected"] = True
                stats["dropbox"]["status"] = "Connected"
            elif 'gmail' in provider:
                stats["gmail"]["connected"] = True
                stats["gmail"]["status"] = "Connected"
                
        return jsonify(stats), 200
    except Exception as e:
        logging.error(f"Error fetching storage stats: {e}")
        return jsonify({"error": "Failed to fetch storage stats"}), 500

@search_bp.route('/file-details', methods=['GET'])
@jwt_required()
def get_file_details():
    user_id = get_jwt_identity()
    file_id = request.args.get('id')
    filepath = request.args.get('filepath')

    if not file_id and not filepath:
        return jsonify({"error": "File ID or filepath is required"}), 400

    query = IndexedFile.query.filter_by(user_id=user_id)
    if file_id:
        file = query.filter_by(id=file_id).first()
    else:
        file = query.filter_by(filepath=filepath).first()

    if not file:
        return jsonify({"error": "File not found"}), 404

    details = file.to_dict()
    details["owner"] = "Me"
    details["size"] = "Unknown"
    details["preview_url"] = None

    try:
        if file.storage_type == "google_drive" and file.cloud_file_id:
            account = CloudStorageAccount.query.filter_by(user_id=user_id, provider="Google Drive").first()
            if account and account.access_token:
                creds = Credentials(
                    token=account.access_token,
                    refresh_token=account.refresh_token,
                    token_uri="https://oauth2.googleapis.com/token",
                    client_id=os.getenv("CLIENT_ID"),
                    client_secret=os.getenv("CLIENT_SECRET")
                )
                drive_service = build("drive", "v3", credentials=creds)
                file_info = drive_service.files().get(
                    fileId=file.cloud_file_id, 
                    fields="thumbnailLink, webContentLink, size, owners"
                ).execute()
                
                details["preview_url"] = file_info.get("thumbnailLink")
                if "size" in file_info:
                    size_bytes = int(file_info['size'])
                    if size_bytes > 1024 * 1024:
                        details["size"] = f"{size_bytes / (1024 * 1024):.1f} MB"
                    else:
                        details["size"] = f"{size_bytes / 1024:.1f} KB"
                if "owners" in file_info and file_info["owners"]:
                    details["owner"] = file_info["owners"][0].get("displayName", "Me")

        elif file.storage_type == "dropbox" and file.cloud_file_id:
            account = CloudStorageAccount.query.filter_by(user_id=user_id, provider="Dropbox").first()
            if account and account.access_token:
                dbx = dropbox.Dropbox(account.access_token)
                try:
                    # The path in Dropbox API usually starts with /
                    db_path = file.filepath.replace('dropbox://', '')
                    res = dbx.files_get_temporary_link(db_path)
                    details["preview_url"] = res.link
                    
                    meta = dbx.files_get_metadata(db_path)
                    if hasattr(meta, 'size'):
                        size_bytes = meta.size
                        if size_bytes > 1024 * 1024:
                            details["size"] = f"{size_bytes / (1024 * 1024):.1f} MB"
                        else:
                            details["size"] = f"{size_bytes / 1024:.1f} KB"
                except Exception as e:
                    logging.warning(f"Failed to get Dropbox preview: {e}")

        elif file.storage_type in ["local", "local_upload"]:
            details["preview_url"] = "local_agent_preview"
            # Attempt to get local file size if it exists on server (local_upload), else leave Unknown
            if file.storage_type == "local_upload" and file.filepath.startswith("upload://"):
                # Can't easily get size without checking disk, leave as Unknown or fetch from DB if we had it
                pass

    except Exception as e:
        logging.error(f"Error fetching file details for preview: {e}")

    return jsonify(details), 200
